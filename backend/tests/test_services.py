import os
import tempfile
import unittest

from pptx import Presentation
from pptx.util import Inches

from app.services.chunker import chunk_text
from app.services.document_parser import parse_document
from app.services.embeddings import generate_embedding
from app.services.vector_store import _collection_name


class TestServiceBehavior(unittest.TestCase):
    # --- Chunking ---
    def test_chunking_empty_document(self):
        result = chunk_text("")
        self.assertEqual(result, [])

    # --- Embeddings ---
    def test_embedding_dimension(self):
        vector = generate_embedding("sample text")
        self.assertEqual(len(vector), 384)

    # --- Document Parsing ---
    def test_parse_document_unsupported_type(self):
        with self.assertRaises(ValueError):
            parse_document("somefile.xyz")

    def test_parse_document_txt_file(self):
        with tempfile.NamedTemporaryFile(mode="w", suffix=".txt", delete=False) as tmp:
            tmp.write("Sample lecture note content for testing.")
            tmp_path = tmp.name
        try:
            pages = parse_document(tmp_path)
            self.assertIsInstance(pages, list)
            self.assertGreaterEqual(len(pages), 1)
            self.assertIn("text", pages[0])
            self.assertIn("page", pages[0])
        finally:
            os.remove(tmp_path)

    def test_parse_document_pptx_file(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Sample Slide"
        body = slide.placeholders[1]
        body.text = "Bullet point one\nBullet point two"

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = tmp.name
        try:
            prs.save(tmp_path)
            pages = parse_document(tmp_path)
            self.assertIsInstance(pages, list)
            self.assertGreaterEqual(len(pages), 1)
            self.assertIn("text", pages[0])
            self.assertIn("page", pages[0])
            self.assertIn("Sample Slide", pages[0]["text"])
        finally:
            os.remove(tmp_path)

    # --- Vector Store ---
    def test_collection_name_format(self):
        name = _collection_name(7)
        self.assertEqual(name, "user_7")


if __name__ == "__main__":
    unittest.main()
