"""
Parse PDF, DOCX, TXT, and PPTX files into plain text with page metadata.
"""
import io
import shutil
import subprocess
from pathlib import Path
from typing import Generator


def parse_document(file_path: str) -> list[dict]:
    """
    Returns a list of page dicts: {"page": int, "text": str}
    """
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return _parse_pdf(file_path)
    elif ext == ".docx":
        return _parse_docx(file_path)
    elif ext == ".txt":
        return _parse_txt(file_path)
    elif ext == ".pptx":
        return _parse_pptx(file_path)
    elif ext == ".ppt":
        return _parse_ppt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def _parse_pdf(file_path: str) -> list[dict]:
    import fitz  # PyMuPDF

    pages = []
    with fitz.open(file_path) as doc:
        for page_num, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if text:
                pages.append({"page": page_num, "text": text})
    return pages


def _parse_docx(file_path: str) -> list[dict]:
    from docx import Document

    doc = Document(file_path)
    # DOCX has no native pages; treat every 30 paragraphs as a "page"
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    pages = []
    page_size = 30
    for i in range(0, max(len(paragraphs), 1), page_size):
        chunk_text = "\n".join(paragraphs[i : i + page_size])
        if chunk_text:
            pages.append({"page": (i // page_size) + 1, "text": chunk_text})
    return pages or [{"page": 1, "text": ""}]


def _parse_txt(file_path: str) -> list[dict]:
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    # Split into ~2000-char "pages"
    page_size = 2000
    pages = []
    for i in range(0, max(len(text), 1), page_size):
        chunk = text[i : i + page_size].strip()
        if chunk:
            pages.append({"page": (i // page_size) + 1, "text": chunk})
    return pages or [{"page": 1, "text": ""}]


def _parse_pptx(file_path: str) -> list[dict]:
    from pptx import Presentation

    presentation = Presentation(file_path)
    pages = []
    for slide_num, slide in enumerate(presentation.slides, start=1):
        text_lines = []
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            text_lines.append(cell_text)
            elif shape.has_text_frame:
                shape_text = shape.text.strip()
                if shape_text:
                    text_lines.append(shape_text)

        if getattr(slide, "has_notes_slide", False) and slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                text_lines.append(notes_text)

        if text_lines:
            pages.append({"page": slide_num, "text": "\n".join(text_lines)})

    return pages or [{"page": 1, "text": ""}]


def _parse_ppt(file_path: str) -> list[dict]:
    if not shutil.which("soffice"):
        raise ValueError(
            "Unsupported file type: .ppt. Please convert to .pptx before uploading."
        )

    output_dir = Path(file_path).parent
    output_path = output_dir / f"{Path(file_path).stem}.pptx"
    try:
        subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to",
                "pptx",
                "--outdir",
                str(output_dir),
                str(file_path),
            ],
            capture_output=True,
            text=True,
            check=True,
        )
        if not output_path.exists():
            raise ValueError("Failed to convert PPT to PPTX for parsing.")
        return _parse_pptx(str(output_path))
    finally:
        try:
            output_path.unlink(missing_ok=True)
        except OSError:
            pass
